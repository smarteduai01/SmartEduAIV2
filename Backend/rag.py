import os
import json
import tempfile
from typing import List, Tuple

import numpy as np
import streamlit as st
import requests
from dotenv import load_dotenv
from pptx import Presentation
import docx2txt
import PyPDF2
from langchain_text_splitters import RecursiveCharacterTextSplitter

# -----------------------------
# ENV + API CONFIG
# -----------------------------
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if not GEMINI_API_KEY:
    raise ValueError("GEMINI_API_KEY is not set in the environment.")

# LLM endpoint (MCQ generation)
GEMINI_ENDPOINT = (
    f"https://generativelanguage.googleapis.com/v1beta/models/"
    f"gemini-2.5-flash:generateContent?key={GEMINI_API_KEY}"
)

# Embedding endpoint (RAG)
GEMINI_EMBEDDING_ENDPOINT = (
    f"https://generativelanguage.googleapis.com/v1beta/models/"
    f"text-embedding-004:batchEmbedContents?key={GEMINI_API_KEY}"
)

# -----------------------------
# FILE TEXT EXTRACTION
# -----------------------------
def extract_text(file_path: str, ext: str) -> str:
    text = ""
    if ext == "pdf":
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    elif ext == "docx":
        text = docx2txt.process(file_path)
    elif ext == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        st.error("Unsupported file type.")
    return text.strip()


# -----------------------------
# GEMINI CALLS: LLM + EMBEDDINGS
# -----------------------------
def call_gemini(prompt: str) -> str | None:
    """Call Gemini 2.5-flash for text generation (MCQs)."""
    headers = {"Content-Type": "application/json"}
    data = {"contents": [{"parts": [{"text": prompt}]}]}
    resp = requests.post(GEMINI_ENDPOINT, headers=headers, json=data)
    resp.raise_for_status()
    json_response = resp.json()
    try:
        return json_response["candidates"][0]["content"]["parts"][0]["text"]
    except KeyError as e:
        st.error(f"Error parsing Gemini response: {e}")
        st.json(json_response)
        return None


def embed_texts(texts: List[str]) -> List[List[float]]:
    """
    Get embeddings for a list of texts using Gemini text-embedding-004.

    Uses batchEmbedContents:
    POST .../models/text-embedding-004:batchEmbedContents
    Request: { "requests": [ { "model": "models/text-embedding-004", "content": {...} }, ... ] }
    Response: { "embeddings": [ { "value": [...] }, ... ] }
    """
    if not texts:
        return []

    headers = {"Content-Type": "application/json"}
    payload = {
        "requests": [
            {
                "model": "models/text-embedding-004",
                "content": {"parts": [{"text": t}]},
            }
            for t in texts
        ]
    }
    resp = requests.post(GEMINI_EMBEDDING_ENDPOINT, headers=headers, json=payload)
    resp.raise_for_status()
    data = resp.json()

    embeddings = []
    try:
        for emb in data["embeddings"]:
            embeddings.append(emb["values"])
    except KeyError as e:
        st.error(f"Error parsing embedding response: {e}")
        st.json(data)
        return []

    return embeddings


# -----------------------------
# SIMPLE IN-MEMORY VECTOR STORE
# -----------------------------
def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    if a.ndim > 1:
        a = a.flatten()
    if b.ndim > 1:
        b = b.flatten()
    denom = (np.linalg.norm(a) * np.linalg.norm(b))
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)


def build_index(chunks: List[str]) -> Tuple[List[str], np.ndarray]:
    """
    Build an in-memory index:
    - chunks: list of chunk texts
    - embeddings: numpy array (num_chunks, dim)
    """
    embeddings_list = embed_texts(chunks)
    if not embeddings_list:
        raise ValueError("Failed to generate embeddings for chunks.")
    embeddings = np.array(embeddings_list, dtype=np.float32)
    return chunks, embeddings


def retrieve_top_k(
    query: str, chunks: List[str], embeddings: np.ndarray, k: int = 5
) -> List[str]:
    """
    Retrieve top-k most similar chunks for a query.
    """
    if not chunks or embeddings is None or len(chunks) == 0:
        return []

    query_emb_list = embed_texts([query])
    if not query_emb_list:
        return []
    query_emb = np.array(query_emb_list[0], dtype=np.float32)

    # Compute cosine similarity with each chunk embedding
    sims = []
    for idx, ch_emb in enumerate(embeddings):
        sim = cosine_similarity(query_emb, ch_emb)
        sims.append((sim, idx))

    # Sort descending by similarity and take top-k
    sims.sort(reverse=True, key=lambda x: x[0])
    top_indices = [idx for _, idx in sims[:k]]

    return [chunks[i] for i in top_indices]


# -----------------------------
# STREAMLIT APP (RAG + MCQ)
# -----------------------------
st.set_page_config(page_title="LLM + RAG MCQ Generator", layout="centered")
st.title("📘 LLM + RAG Based MCQ Generator")
st.write(
    "Upload a PDF, DOCX, or PPTX — I’ll extract the key concepts, build a RAG index, "
    "and generate simple, meaningful MCQs in JSON format."
)

uploaded_file = st.file_uploader("Upload your file:", type=["pdf", "docx", "pptx"])

if uploaded_file:
    ext = uploaded_file.name.split(".")[-1].lower()
    tmp_file_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix="." + ext) as tmp:
            tmp.write(uploaded_file.read())
            tmp_file_path = tmp.name

        text = extract_text(tmp_file_path, ext)
    finally:
        if tmp_file_path and os.path.exists(tmp_file_path):
            os.unlink(tmp_file_path)

    if text:
        st.success("✅ File processed successfully!")

        # Preview
        st.write("### Extracted Preview:")
        st.text_area(
            "Document Preview",
            text[:1000] + ("..." if len(text) > 1000 else ""),
            height=200,
        )

        # Build chunks + vector index (RAG) once per upload
        with st.spinner("🔧 Building RAG index (chunking + embeddings)..."):
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1500,
                chunk_overlap=200,
            )
            chunks = splitter.split_text(text)

            # Build in-memory index
            try:
                indexed_chunks, chunk_embeddings = build_index(chunks)
                st.session_state["rag_chunks"] = indexed_chunks
                st.session_state["rag_embeddings"] = chunk_embeddings
                st.success("📚 RAG index ready! You can now generate MCQs.")
            except Exception as e:
                st.error(f"Failed to build RAG index: {e}")
                st.stop()

        # Optional: let the user specify how many questions or what focus
        st.write("### MCQ Generation Settings")
        num_questions = st.slider(
            "Number of MCQs to generate", min_value=3, max_value=20, value=8, step=1
        )
        user_focus = st.text_input(
            "Optional: Focus area (e.g., 'Unit 2 only', 'definitions', 'algorithms').",
            "",
        )

        if st.button("Generate MCQs"):
            with st.spinner("🧠 Retrieving relevant content and generating MCQs..."):
                rag_chunks = st.session_state.get("rag_chunks", [])
                rag_embeddings = st.session_state.get("rag_embeddings", None)

                if not rag_chunks or rag_embeddings is None:
                    st.error("RAG index not found. Please re-upload the file.")
                    st.stop()

                # Build a retrieval query for RAG
                if user_focus.strip():
                    query = (
                        f"Key concepts related to: {user_focus.strip()}. "
                        f"Use this to generate {num_questions} simple, meaningful MCQs."
                    )
                else:
                    query = (
                        f"Key concepts that are most important in this document "
                        f"for generating about {num_questions} simple, meaningful MCQs."
                    )

                # Retrieve top-k chunks (RAG core step)
                retrieved_chunks = retrieve_top_k(
                    query, rag_chunks, rag_embeddings, k=5
                )

                if not retrieved_chunks:
                    st.error("Could not retrieve relevant chunks for MCQ generation.")
                    st.stop()

                context_text = "\n\n".join(retrieved_chunks)

                # Prompt with retrieved context only (RAG)
                prompt = f"""
You are an expert educational AI system. Your task is to generate multiple-choice questions (MCQs)
based on the most important and easy-to-understand themes in the given content.

### Retrieved Context:
{context_text}

### Instructions:
- Generate **{num_questions}** simple, meaningful MCQs that test understanding.
- Each question must have exactly 4 options and one correct answer.
- Assign a difficulty level to each question based on its conceptual depth:
  - "Easy" → simple recall or fact-based
  - "Medium" → involves understanding or application
  - "Hard" → requires analysis or reasoning
- Use information **only** from the retrieved context above.
- Return the output strictly in this JSON format (with many entries, one per question):

{{
  "Question text 1": {{
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_option": "Correct Option Text",
    "difficulty": "Easy | Medium | Hard"
  }},
  "Question text 2": {{
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_option": "Correct Option Text",
    "difficulty": "Easy | Medium | Hard"
  }}
  ...
}}

Do not include any explanations or text outside JSON.
"""

                try:
                    response_text = call_gemini(prompt)
                    if response_text:
                        # Strip ```json fences if model adds them
                        cleaned = response_text.strip()
                        if cleaned.startswith("```json"):
                            cleaned = cleaned[7:]
                        if cleaned.endswith("```"):
                            cleaned = cleaned[:-3]
                        cleaned = cleaned.strip()

                        mcq_json = json.loads(cleaned)
                        st.success("✅ MCQs generated successfully!")
                        st.json(mcq_json)

                        json_str = json.dumps(mcq_json, indent=4)
                        st.download_button(
                            label="📥 Download JSON",
                            data=json_str,
                            file_name="generated_mcqs.json",
                            mime="application/json",
                        )
                    else:
                        st.error("No response received from Gemini.")
                except json.JSONDecodeError:
                    st.error("❌ LLM response was not valid JSON. Please retry.")
                    st.text("Raw response from LLM:")
                    st.text(response_text)
                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")
                    st.text("Raw response from LLM (if available):")
                    st.text(response_text if "response_text" in locals() else "No response captured.")

    else:
        st.error(
            "❌ Could not extract any text from the uploaded file. "
            "It might be empty, scanned as an image, or corrupted."
        )
